import GPT
from pptx import Presentation
from pptx.util import *
def ppt_detail(title, wordcount, language, tone, decision):
        if decision == 'a':
                global x
                x = Presentation()
                # layout
                Layout = x.slide_layouts[1] 
                # # slides
                first_slide = x.slides.add_slide(Layout)
                first_slide.shapes.title.text = GPT.catch_title(title,language)
                #import response form GPT
                response = GPT.title_overview(title,wordcount,language,tone)
                first_slide.placeholders[1].text = response
        elif decision == 'b':
                global template_user
                template_user = Presentation("D:\\major project\\ppt_template_user\\NCR_CU_PPT_Template_2019-2020.pptx")
                # layout
                Layout = template_user.slide_layouts[1]
                # # slides
                first_slide = template_user.slides.add_slide(Layout)
                first_slide.shapes.title.text = GPT.catch_title(title, language)
                #import response form GPT
                response = GPT.title_overview(title, wordcount, language, tone)
                first_slide.placeholders[0].text = response
        else:
                print("Invalid input")      
                
                
def ppt_keywords(key, keywords, title, wordcount, language, tone, decision):
        if decision == 'a':
                #for each keyword make a slide and wordcount is the number of words in the slide
                for i in range(key):
                        slide = x.slides.add_slide(x.slide_layouts[1])
                        # slide.shapes.title.text = GPT.translate_title(keywords[i],language)
                        slide.shapes.title.text = keywords[i]
                        response = GPT.keyword_overview(keywords[i],wordcount,language,tone)
                        slide.placeholders[1].text = response
                x.save('testing.pptx')    
        
        elif decision == 'b':
               #for each keyword make a slide and wordcount is the number of words in the slide
                for i in range(key):
                        slide = template_user.slides.add_slide(template_user.slide_layouts[1])
                        slide.shapes.title.text = GPT.translate_title(keywords[i], language)
                        response = GPT.keyword_overview(keywords[i], wordcount, language, tone)
                        slide.placeholders[1].text = response
                template_user.save('testing.pptx')
         
        else:
                print("Invalid input")
        